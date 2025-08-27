
import streamlit as st
import fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont
import numpy as np
from streamlit_drawable_canvas import st_canvas
from pptx import Presentation
from pptx.util import Inches, Pt
import io, zipfile, os

st.set_page_config(page_title="PDF Editor (Streamlit)", layout="wide")

# -----------------------------
# Utilities
# -----------------------------

def pdf_bytes_to_images(pdf_bytes: bytes, target_width: int = 1200):
    """Render each page of a PDF (bytes) to a PIL Image with approximately target_width (px)."""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    images = []
    for i, page in enumerate(doc):
        # Calculate zoom factor based on desired width
        page_width = page.rect.width  # in points
        zoom = target_width / page_width
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=True)
        mode = "RGBA" if pix.alpha else "RGB"
        img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
        images.append(img)
    doc.close()
    return images

def thumbnail(img: Image.Image, max_w=220):
    w, h = img.size
    if w > max_w:
        ratio = max_w / w
        img = img.resize((int(w * ratio), int(h * ratio)))
    return img

def hex_to_rgba(s: str, alpha: int = 255):
    s = s.lstrip("#")
    if len(s) == 6:
        r = int(s[0:2], 16)
        g = int(s[2:4], 16)
        b = int(s[4:6], 16)
        return (r, g, b, alpha)
    return (0, 0, 0, alpha)

def draw_texts_on_image(img: Image.Image, texts: list, default_font_path: str | None):
    """Draw a list of text annotations onto img.
    texts item: {"text": str, "x": int, "y": int, "size": int, "color": "#RRGGBB", "font_path": str|None}"""
    if not texts:
        return img
    out = img.copy()
    draw = ImageDraw.Draw(out)
    for t in texts:
        content = t.get("text","")
        x = int(t.get("x",0))
        y = int(t.get("y",0))
        size = int(t.get("size",24))
        color_hex = t.get("color","#000000")
        font_path = t.get("font_path") or default_font_path
        try:
            if font_path:
                font = ImageFont.truetype(font_path, size=size)
            else:
                font = ImageFont.load_default()
        except Exception:
            font = ImageFont.load_default()
        draw.text((x,y), content, fill=hex_to_rgba(color_hex, 255), font=font)
    return out

def compose_final_page(base_img: Image.Image, canvas_img: np.ndarray | None, text_annos: list, font_path: str | None):
    """Merge the canvas result (already contains base) + additional texts. If canvas_img is None, use base_img."""
    if canvas_img is not None:
        pil_canvas = Image.fromarray(canvas_img.astype("uint8"), "RGBA")
        merged = pil_canvas
    else:
        merged = base_img.copy()
    merged = draw_texts_on_image(merged, text_annos, font_path)
    return merged

def save_as_pdf(images: list[Image.Image]) -> bytes:
    if not images:
        return b""
    buf = io.BytesIO()
    # Convert all to RGB for PDF (avoid transparency issues)
    rgb_imgs = []
    for im in images:
        if im.mode != "RGB":
            rgb_imgs.append(im.convert("RGB"))
        else:
            rgb_imgs.append(im)
    rgb_imgs[0].save(buf, format="PDF", save_all=True, append_images=rgb_imgs[1:])
    buf.seek(0)
    return buf.read()

def save_as_images_zip(images: list[Image.Image], fmt: str = "PNG") -> bytes:
    """Return a ZIP containing page_001.(fmt), ..."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, im in enumerate(images, start=1):
            b = io.BytesIO()
            if fmt.upper() == "JPG" or fmt.upper() == "JPEG":
                # Ensure no alpha for JPEG
                to_save = im.convert("RGB")
                ext = "jpg"
            else:
                to_save = im
                ext = "png"
            to_save.save(b, format=fmt.upper())
            zf.writestr(f"page_{i:03d}.{ext}", b.getvalue())
    buf.seek(0)
    return buf.read()

def save_as_pptx(images: list[Image.Image]) -> bytes:
    prs = Presentation()
    # optional: set slide size to 16:9 for consistency
    prs.slide_width = Inches(13.333)   # 1280px @96dpi approx
    prs.slide_height = Inches(7.5)     # 720px  @96dpi approx
    blank_layout = prs.slide_layouts[6]  # blank
    for im in images:
        slide = prs.slides.add_slide(blank_layout)
        # Save image to bytes
        b = io.BytesIO()
        # Convert to RGB if needed
        if im.mode != "RGB":
            im_rgb = im.convert("RGB")
        else:
            im_rgb = im
        im_rgb.save(b, format="PNG")
        b.seek(0)
        pic = slide.shapes.add_picture(b, left=Inches(0), top=Inches(0))
        # Fit to slide while preserving aspect ratio
        slide_w, slide_h = prs.slide_width, prs.slide_height
        img_w, img_h = pic.width, pic.height
        scale = min(slide_w/img_w, slide_h/img_h)
        pic.width = int(img_w * scale)
        pic.height = int(img_h * scale)
        # Center
        pic.left = int((slide_w - pic.width) / 2)
        pic.top = int((slide_h - pic.height) / 2)
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()

def ensure_session():
    st.session_state.setdefault("pdf_bytes", None)
    st.session_state.setdefault("images", [])                # base page images for editing (at chosen width)
    st.session_state.setdefault("edited_canvas", {})         # page_idx -> np.ndarray (RGBA) from st_canvas.image_data
    st.session_state.setdefault("text_annos", {})            # page_idx -> list of text dicts
    st.session_state.setdefault("keep_flags", [])            # per page True/False
    st.session_state.setdefault("order", [])                 # list of ints 0..n-1 indicating current order
    st.session_state.setdefault("font_file_bytes", None)     # uploaded TTF font
    st.session_state.setdefault("font_path_tmp", None)       # temp saved TTF path

ensure_session()

st.title("📄 Streamlit PDF 편집기")

with st.sidebar:
    st.header("⚙️ 설정")
    work_width = st.slider("작업 해상도(가로 px)", min_value=800, max_value=2000, value=1200, step=100,
                           help="페이지를 이 해상도로 렌더링하여 편집합니다. 값이 높을수록 품질↑/속도↓")
    uploaded_font = st.file_uploader("한글 폰트 업로드(.ttf 권장)", type=["ttf", "otf"], accept_multiple_files=False,
                                     help="텍스트 입력 시 적용. 미업로드 시 기본 폰트 사용(한글 일부 미지원 가능).")
    if uploaded_font is not None:
        st.session_state["font_file_bytes"] = uploaded_font.read()
        # Save temp font file
        tmp_path = os.path.join(st.experimental_get_query_params().get("font_dir",[os.getcwd()])[0], "uploaded_font.ttf")
        try:
            with open(tmp_path, "wb") as f:
                f.write(st.session_state["font_file_bytes"])
            st.session_state["font_path_tmp"] = tmp_path
            st.success("폰트 업로드 완료!")
        except Exception as e:
            st.warning(f"폰트를 저장하지 못했습니다: {e}")

st.markdown("**기능:** PDF 업로드 → 페이지 순서 변경/삭제 → 도형/펜으로 그리기 → 텍스트 추가 → 저장( PDF / PNG(zip) / JPG(zip) / PPTX )")

# -----------------------------
# Step 1. Upload
# -----------------------------
file = st.file_uploader("PDF 파일 업로드", type=["pdf"])
if file is not None:
    st.session_state["pdf_bytes"] = file.read()
    # Reset everything when a new file is uploaded
    st.session_state["images"] = pdf_bytes_to_images(st.session_state["pdf_bytes"], target_width=work_width)
    n = len(st.session_state["images"])
    st.session_state["edited_canvas"] = {}
    st.session_state["text_annos"] = {i: [] for i in range(n)}
    st.session_state["keep_flags"] = [True]*n
    st.session_state["order"] = list(range(n))

if not st.session_state["images"]:
    st.info("왼쪽 상단에서 PDF를 업로드해 주세요.")
    st.stop()

# -----------------------------
# Step 2. Manage pages (order/delete)
# -----------------------------
st.subheader("1) 페이지 관리 – 순서 변경 & 삭제")
imgs = st.session_state["images"]
n = len(imgs)

cols = st.columns(4)
for i, img in enumerate(imgs):
    col = cols[i % 4]
    with col:
        st.image(thumbnail(img), caption=f"페이지 {i+1}", use_column_width=True)
        st.session_state["keep_flags"][i] = st.checkbox(f"유지(삭제 해제)", value=st.session_state["keep_flags"][i], key=f"keep_{i}")
        # Order input as 1..n
        new_pos = st.number_input(f"순서", min_value=1, max_value=n, value=st.session_state["order"].index(i)+1, key=f"order_input_{i}")
        # Store desired order indirectly via a separate structure
        # We'll reconstruct after the loop

# Reconstruct order based on user's numbers (resolve conflicts by original index)
desired = [(st.session_state[f"order_input_{i}"], i) for i in range(n)]
# sort by (desired_position, original_index)
desired_sorted = sorted(desired, key=lambda x: (x[0], x[1]))
new_order = [idx for _, idx in desired_sorted]
st.session_state["order"] = new_order

st.caption("👉 팁: '순서' 숫자를 바꾼 뒤 아래 편집 또는 저장 단계로 바로 진행하면 반영됩니다.")

# -----------------------------
# Step 3. Page editor (draw shapes + add texts)
# -----------------------------
st.subheader("2) 페이지 편집 – 도형/펜 & 텍스트")

# Choose a page to edit (after ordering & deletion flags but before applying)
# We'll show selection based on the **current** order, but indicate original index
valid_indices = [i for i in st.session_state["order"] if st.session_state["keep_flags"][i]]
if not valid_indices:
    st.warning("모든 페이지가 삭제로 설정되어 있습니다. 최소 1페이지는 유지해 주세요.")
    st.stop()

page_choice = st.selectbox(
    "편집할 페이지 선택 (현재 순서 기준)",
    options=valid_indices,
    format_func=lambda i: f"페이지 {i+1} (원본 번호)"
)

page_img = st.session_state["images"][page_choice]

left, right = st.columns([3,2], vertical_alignment="top")

with left:
    st.markdown("**도형/펜 그리기**")
    drawing_mode = st.selectbox("드로잉 모드", ["freedraw", "line", "rect", "circle", "transform"], index=0,
                                help="텍스트는 우측 '텍스트 추가' 섹션에서 별도 입력")
    stroke_width = st.slider("선 굵기", 1, 25, 3)
    stroke_color = st.color_picker("선 색상", "#000000")
    fill_color = st.color_picker("채우기 색상(도형용)", "#000000")

    # Create canvas; background is the page image
    canvas_result = st_canvas(
        fill_color=fill_color + "88",  # plus alpha channel
        stroke_width=stroke_width,
        stroke_color=stroke_color,
        background_image=page_img,
        update_streamlit=True,
        height=page_img.height,
        width=page_img.width,
        drawing_mode=drawing_mode,
        initial_drawing=None,
        display_toolbar=True,
        key=f"canvas_{page_choice}"
    )

with right:
    st.markdown("**텍스트 추가**")
    # Prepare per-page text list
    if page_choice not in st.session_state["text_annos"]:
        st.session_state["text_annos"][page_choice] = []

    # Add new text entry
    with st.expander("텍스트 박스 추가/관리", expanded=True):
        add_col1, add_col2 = st.columns(2)
        with add_col1:
            new_text = st.text_input("텍스트 내용", key=f"new_text_{page_choice}")
            txt_color = st.color_picker("글자 색상", "#000000", key=f"txt_color_{page_choice}")
            txt_size = st.slider("글자 크기", 8, 128, 28, key=f"txt_size_{page_choice}")
        with add_col2:
            # Place using sliders (px). Tip: 이미지를 보며 대략 위치를 지정
            x = st.number_input("X 위치(px)", min_value=0, max_value=page_img.width, value=50, step=5, key=f"x_{page_choice}")
            y = st.number_input("Y 위치(px)", min_value=0, max_value=page_img.height, value=50, step=5, key=f"y_{page_choice}")
            font_path = st.session_state.get("font_path_tmp")

        add_btn = st.button("이 텍스트 추가", key=f"add_text_btn_{page_choice}")
        if add_btn and new_text.strip():
            st.session_state["text_annos"][page_choice].append({
                "text": new_text.strip(),
                "x": int(x),
                "y": int(y),
                "size": int(txt_size),
                "color": txt_color,
                "font_path": font_path
            })
            st.success("텍스트를 추가했습니다.")

        # List existing texts
        if st.session_state["text_annos"][page_choice]:
            st.write("현재 텍스트들:")
            to_delete = st.multiselect(
                "삭제할 텍스트 선택",
                options=list(range(len(st.session_state["text_annos"][page_choice]))),
                format_func=lambda idx: f"[{idx}] '{st.session_state['text_annos'][page_choice][idx]['text']}' @({st.session_state['text_annos'][page_choice][idx]['x']},{st.session_state['text_annos'][page_choice][idx]['y']}) size={st.session_state['text_annos'][page_choice][idx]['size']}"
            )
            if st.button("선택 텍스트 삭제", key=f"del_text_btn_{page_choice}") and to_delete:
                # Remove in reverse order to keep indices stable
                for idx in sorted(to_delete, reverse=True):
                    st.session_state["text_annos"][page_choice].pop(idx)
                st.success("선택한 텍스트를 삭제했습니다.")

    st.markdown("---")
    st.markdown("**미리보기/적용**")
    apply_btn = st.button("현재 페이지 적용(그림 + 텍스트) 저장", key=f"apply_btn_{page_choice}")
    if apply_btn:
        canvas_img = canvas_result.image_data if canvas_result is not None else None
        if canvas_img is None:
            # If nothing drawn yet, create from base to allow text only
            base_for_canvas = page_img.copy()
            canvas_img = np.array(base_for_canvas.convert("RGBA"))
        final = compose_final_page(page_img, canvas_img, st.session_state["text_annos"].get(page_choice, []),
                                   st.session_state.get("font_path_tmp"))
        st.session_state["edited_canvas"][page_choice] = np.array(final.convert("RGBA"))
        st.success("이 페이지의 편집 내용을 저장했습니다.")

# Show a preview of the currently selected page (edited if any)
st.markdown("**현재 선택 페이지 미리보기**")
preview_img = None
if page_choice in st.session_state["edited_canvas"]:
    preview_img = Image.fromarray(st.session_state["edited_canvas"][page_choice], "RGBA")
else:
    if canvas_result and canvas_result.image_data is not None:
        preview_img = Image.fromarray(canvas_result.image_data.astype("uint8"), "RGBA")
    else:
        preview_img = page_img
st.image(preview_img, use_column_width=True)

# -----------------------------
# Step 4. Export
# -----------------------------
st.subheader("3) 저장하기")

fmt = st.selectbox("저장 형식", ["PDF", "PNG(zip)", "JPG(zip)", "PPTX"], index=0)

if st.button("저장 파일 만들기"):
    # Build final ordered + kept pages as PIL images
    final_pages = []
    for i in st.session_state["order"]:
        if not st.session_state["keep_flags"][i]:
            continue
        if i in st.session_state["edited_canvas"]:
            final_img = Image.fromarray(st.session_state["edited_canvas"][i], "RGBA")
        else:
            # If user didn't click "적용", fall back to base or current canvas state
            if st.session_state.get(f"canvas_{i}") and st.session_state[f"canvas_{i}"].image_data is not None:
                final_img = Image.fromarray(st.session_state[f"canvas_{i}"].image_data.astype("uint8"), "RGBA")
            else:
                final_img = st.session_state["images"][i]
        final_pages.append(final_img)

    if not final_pages:
        st.warning("저장할 페이지가 없습니다. 최소 1페이지 이상 유지해 주세요.")
    else:
        if fmt == "PDF":
            data = save_as_pdf(final_pages)
            st.download_button("📥 PDF 다운로드", data=data, file_name="edited.pdf", mime="application/pdf")
        elif fmt == "PNG(zip)":
            data = save_as_images_zip(final_pages, fmt="PNG")
            st.download_button("📥 PNG ZIP 다운로드", data=data, file_name="pages_png.zip", mime="application/zip")
        elif fmt == "JPG(zip)":
            data = save_as_images_zip(final_pages, fmt="JPEG")
            st.download_button("📥 JPG ZIP 다운로드", data=data, file_name="pages_jpg.zip", mime="application/zip")
        else:  # PPTX
            data = save_as_pptx(final_pages)
            st.download_button("📥 PPTX 다운로드", data=data, file_name="slides.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

st.caption("※ 성능 팁: 페이지가 크거나 많다면 좌측 '작업 해상도'를 낮춰주세요. 텍스트 한글 폰트는 사이드바에서 업로드 가능합니다.")
