
import streamlit as st
import fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont
import numpy as np
from streamlit_drawable_canvas import st_canvas
from pptx import Presentation
from pptx.util import Inches
import io, zipfile, os, base64

st.set_page_config(page_title="PDF Editor (Streamlit) – bg fix", layout="wide")

def pdf_bytes_to_images(pdf_bytes: bytes, target_width: int = 1200):
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    images = []
    for page in doc:
        page_width = page.rect.width
        zoom = target_width / page_width
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=True)
        mode = "RGBA" if pix.alpha else "RGB"
        img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
        images.append(img)
    doc.close()
    return images

def thumbnail(img: Image.Image, max_w=220):
    w, h = img.size
    if w > max_w:
        ratio = max_w / w
        img = img.resize((int(w * ratio), int(h * ratio)))
    return img

def hex_with_alpha(hex_color: str, alpha_0_255: int) -> str:
    h = hex_color.lstrip("#")
    return f"#{h}{alpha_0_255:02x}"

def hex_to_rgba_tuple(hex_color: str, alpha_0_255: int):
    h = hex_color.lstrip("#")
    r = int(h[0:2], 16); g = int(h[2:4], 16); b = int(h[4:6], 16)
    return (r, g, b, max(0, min(255, alpha_0_255)))

def draw_texts_on_image(img: Image.Image, texts: list, default_font_path: str | None):
    if not texts: return img
    out = img.copy()
    draw = ImageDraw.Draw(out)
    for t in texts:
        content = t.get("text","")
        x = int(t.get("x",0)); y = int(t.get("y",0))
        size = int(t.get("size",24))
        color_hex = t.get("color","#000000")
        font_path = t.get("font_path") or default_font_path
        try:
            font = ImageFont.truetype(font_path, size=size) if font_path else ImageFont.load_default()
        except Exception:
            font = ImageFont.load_default()
        rgba = hex_to_rgba_tuple(color_hex, 255)
        draw.text((x,y), content, fill=rgba, font=font)
    return out

def compose_final_page(base_img, canvas_img, text_annos, font_path):
    if canvas_img is not None:
        pil_canvas = Image.fromarray(canvas_img.astype("uint8"), "RGBA")
        merged = pil_canvas
    else:
        merged = base_img.copy()
    merged = draw_texts_on_image(merged, text_annos, font_path)
    return merged

def save_as_pdf(images):
    if not images: return b""
    buf = io.BytesIO()
    rgb_imgs = [(im.convert("RGB") if im.mode != "RGB" else im) for im in images]
    rgb_imgs[0].save(buf, format="PDF", save_all=True, append_images=rgb_imgs[1:])
    buf.seek(0); return buf.read()

def save_as_images_zip(images, fmt="PNG"):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, im in enumerate(images, start=1):
            b = io.BytesIO()
            if fmt.upper() in ("JPG","JPEG"):
                to_save = im.convert("RGB"); ext="jpg"
            else:
                to_save = im; ext="png"
            to_save.save(b, format=fmt.upper())
            zf.writestr(f"page_{i:03d}.{ext}", b.getvalue())
    buf.seek(0); return buf.read()

def save_as_pptx(images):
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for im in images:
        slide = prs.slides.add_slide(blank)
        b = io.BytesIO()
        (im.convert("RGB") if im.mode != "RGB" else im).save(b, format="PNG"); b.seek(0)
        pic = slide.shapes.add_picture(b, left=Inches(0), top=Inches(0))
        slide_w, slide_h = prs.slide_width, prs.slide_height
        img_w, img_h = pic.width, pic.height
        scale = min(slide_w/img_w, slide_h/img_h)
        pic.width = int(img_w * scale); pic.height = int(img_h * scale)
        pic.left = int((slide_w - pic.width)/2); pic.top = int((slide_h - pic.height)/2)
    out = io.BytesIO(); prs.save(out); out.seek(0); return out.read()

def ensure_session():
    st.session_state.setdefault("pdf_bytes", None)
    st.session_state.setdefault("images", [])
    st.session_state.setdefault("edited_canvas", {})
    st.session_state.setdefault("text_annos", {})
    st.session_state.setdefault("keep_flags", [])
    st.session_state.setdefault("order", [])
    st.session_state.setdefault("font_file_bytes", None)
    st.session_state.setdefault("font_path_tmp", None)

ensure_session()

st.title("📄 Streamlit PDF 편집기 (배경버그 우회판)")

with st.sidebar:
    st.header("⚙️ 설정")
    work_width = st.slider("작업 해상도(가로 px)", 800, 2000, 1200, 100)
    uploaded_font = st.file_uploader("한글 폰트 업로드(.ttf/.otf)", type=["ttf","otf"], accept_multiple_files=False)
    if uploaded_font is not None:
        st.session_state["font_file_bytes"] = uploaded_font.read()
        tmp_path = os.path.join(os.getcwd(), "uploaded_font.ttf")
        try:
            with open(tmp_path, "wb") as f: f.write(st.session_state["font_file_bytes"])
            st.session_state["font_path_tmp"] = tmp_path
            st.success("폰트 업로드 완료")
        except Exception as e:
            st.warning(f"폰트 저장 실패: {e}")

st.markdown("**PDF 업로드 → 순서/삭제 → 도형/펜/텍스트(투명도) → 저장(PDF/PNG/JPG/PPTX)**")

# ---- Upload ----
file = st.file_uploader("PDF 파일 업로드", type=["pdf"])
if file is not None:
    st.session_state["pdf_bytes"] = file.read()
    st.session_state["images"] = pdf_bytes_to_images(st.session_state["pdf_bytes"], target_width=work_width)
    n = len(st.session_state["images"])
    st.session_state["edited_canvas"] = {}
    st.session_state["text_annos"] = {i: [] for i in range(n)}
    st.session_state["keep_flags"] = [True]*n
    st.session_state["order"] = list(range(n))

if not st.session_state["images"]:
    st.info("왼쪽에서 PDF를 업로드해 주세요.")
    st.stop()

# ---- Manage pages ----
st.subheader("1) 페이지 관리 – 순서 변경 & 삭제")
imgs = st.session_state["images"]; n = len(imgs)
cols = st.columns(4)
for i, img in enumerate(imgs):
    col = cols[i % 4]
    with col:
        st.image(thumbnail(img), caption=f"페이지 {i+1}", use_container_width=True)
        st.session_state["keep_flags"][i] = st.checkbox("유지(삭제 해제)", value=st.session_state["keep_flags"][i], key=f"keep_{i}")
        new_pos = st.number_input("순서", 1, n, value=st.session_state["order"].index(i)+1, key=f"order_input_{i}")

desired = [(st.session_state[f"order_input_{i}"], i) for i in range(n)]
desired_sorted = sorted(desired, key=lambda x: (x[0], x[1]))
st.session_state["order"] = [idx for _, idx in desired_sorted]
st.caption("순서를 바꾼 뒤 아래 편집/저장 단계로 진행하면 반영됩니다.")

# ---- Editor ----
st.subheader("2) 페이지 편집 – 도형/펜 & 텍스트")
valid_indices = [i for i in st.session_state["order"] if st.session_state["keep_flags"][i]]
if not valid_indices:
    st.warning("모든 페이지를 삭제하셨습니다. 최소 1페이지를 유지하세요.")
    st.stop()

page_choice = st.selectbox("편집할 페이지 선택(원본 번호)", options=valid_indices, format_func=lambda i: f"페이지 {i+1}")
page_img = st.session_state["images"][page_choice]

left, right = st.columns([3,2], vertical_alignment="top")
with left:
    st.markdown("**도형/펜 그리기**")
    drawing_mode = st.selectbox("드로잉 모드", ["freedraw","line","rect","circle","transform"], index=0)
    stroke_width = st.slider("선 굵기", 1, 25, 3)
    stroke_color = st.color_picker("선 색상", "#000000")
    stroke_alpha_pct = st.slider("선 투명도(%)", 0, 100, 100)
    fill_color = st.color_picker("채우기 색상(도형용)", "#000000")
    fill_alpha_pct = st.slider("채우기 투명도(%)", 0, 100, 50)

    stroke_alpha = int(255 * (stroke_alpha_pct/100.0))
    fill_alpha = int(255 * (fill_alpha_pct/100.0))
    stroke_color_rgba = hex_with_alpha(stroke_color, stroke_alpha)
    fill_color_rgba = hex_with_alpha(fill_color, fill_alpha)

    # Build Fabric.js initial state that embeds the page image as base layer (unselectable)
    buf = io.BytesIO()
    page_img.save(buf, format="PNG"); b64 = base64.b64encode(buf.getvalue()).decode("ascii")
    data_url = f"data:image/png;base64,{b64}"
    initial_state = {
        "version": "4.4.0",
        "objects": [
            {
                "type": "image",
                "version": "4.4.0",
                "originX": "left",
                "originY": "top",
                "left": 0, "top": 0,
                "scaleX": 1, "scaleY": 1,
                "angle": 0,
                "opacity": 1,
                "selectable": False,
                "evented": False,
                "hasControls": False,
                "hasBorders": False,
                "src": data_url
            }
        ]
    }

    canvas_result = st_canvas(
        fill_color=fill_color_rgba,
        stroke_width=stroke_width,
        stroke_color=stroke_color_rgba,
        background_color="#00000000",  # transparent
        update_streamlit=True,
        height=page_img.height,
        width=page_img.width,
        drawing_mode=drawing_mode,
        initial_drawing=initial_state,   # <-- embed page image here
        display_toolbar=True,
        key=f"canvas_{page_choice}"
    )

with right:
    st.markdown("**텍스트 추가**")
    if page_choice not in st.session_state["text_annos"]:
        st.session_state["text_annos"][page_choice] = []
    with st.expander("텍스트 박스 추가/관리", expanded=True):
        add_col1, add_col2 = st.columns(2)
        with add_col1:
            new_text = st.text_input("텍스트 내용", key=f"new_text_{page_choice}")
            txt_color = st.color_picker("글자 색상", "#000000", key=f"txt_color_{page_choice}")
            txt_size = st.slider("글자 크기", 8, 128, 28, key=f"txt_size_{page_choice}")
        with add_col2:
            x = st.number_input("X 위치(px)", 0, page_img.width, 50, 5, key=f"x_{page_choice}")
            y = st.number_input("Y 위치(px)", 0, page_img.height, 50, 5, key=f"y_{page_choice}")
            font_path = st.session_state.get("font_path_tmp")
        if st.button("이 텍스트 추가", key=f"add_text_btn_{page_choice}") and new_text.strip():
            st.session_state["text_annos"][page_choice].append({
                "text": new_text.strip(), "x": int(x), "y": int(y),
                "size": int(txt_size), "color": txt_color, "font_path": font_path
            })
            st.success("텍스트를 추가했습니다.")
        if st.session_state["text_annos"][page_choice]:
            st.write("현재 텍스트들:")
            to_delete = st.multiselect(
                "삭제할 텍스트 선택",
                options=list(range(len(st.session_state["text_annos"][page_choice]))),
                format_func=lambda idx: f"[{idx}] '{st.session_state['text_annos'][page_choice][idx]['text']}'@({st.session_state['text_annos'][page_choice][idx]['x']},{st.session_state['text_annos'][page_choice][idx]['y']}) size={st.session_state['text_annos'][page_choice][idx]['size']}"
            )
            if st.button("선택 텍스트 삭제", key=f"del_text_btn_{page_choice}") and to_delete:
                for idx in sorted(to_delete, reverse=True):
                    st.session_state["text_annos"][page_choice].pop(idx)
                st.success("삭제했습니다.")

    st.markdown("---")
    apply_btn = st.button("현재 페이지 적용(그림 + 텍스트) 저장", key=f"apply_btn_{page_choice}")
    if apply_btn:
        canvas_img = canvas_result.image_data if canvas_result is not None else None
        if canvas_img is None:
            # fallback: compose with base only
            canvas_img = np.array(page_img.convert("RGBA"))
        final = compose_final_page(page_img, canvas_img, st.session_state["text_annos"].get(page_choice, []),
                                   st.session_state.get("font_path_tmp"))
        st.session_state["edited_canvas"][page_choice] = np.array(final.convert("RGBA"))
        st.success("이 페이지의 편집 내용을 저장했습니다.")

st.markdown("**현재 선택 페이지 미리보기**")
if page_choice in st.session_state["edited_canvas"]:
    preview_img = Image.fromarray(st.session_state["edited_canvas"][page_choice], "RGBA")
elif canvas_result and canvas_result.image_data is not None:
    preview_img = Image.fromarray(canvas_result.image_data.astype("uint8"), "RGBA")
else:
    preview_img = page_img
st.image(preview_img, use_container_width=True)

# ---- Export ----
st.subheader("3) 저장하기")
fmt = st.selectbox("저장 형식", ["PDF","PNG(zip)","JPG(zip)","PPTX"], index=0)

if st.button("저장 파일 만들기"):
    final_pages = []
    for i in st.session_state["order"]:
        if not st.session_state["keep_flags"][i]: continue
        if i in st.session_state["edited_canvas"]:
            final_img = Image.fromarray(st.session_state["edited_canvas"][i], "RGBA")
        else:
            key = f"canvas_{i}"
            if st.session_state.get(key) and getattr(st.session_state[key], "image_data", None) is not None:
                final_img = Image.fromarray(st.session_state[key].image_data.astype("uint8"), "RGBA")
            else:
                final_img = st.session_state["images"][i]
        final_pages.append(final_img)

    if not final_pages:
        st.warning("저장할 페이지가 없습니다.")
    else:
        if fmt == "PDF":
            data = save_as_pdf(final_pages)
            st.download_button("📥 PDF 다운로드", data=data, file_name="edited.pdf", mime="application/pdf")
        elif fmt == "PNG(zip)":
            data = save_as_images_zip(final_pages, fmt="PNG")
            st.download_button("📥 PNG ZIP 다운로드", data=data, file_name="pages_png.zip", mime="application/zip")
        elif fmt == "JPG(zip)":
            data = save_as_images_zip(final_pages, fmt="JPEG")
            st.download_button("📥 JPG ZIP 다운로드", data=data, file_name="pages_jpg.zip", mime="application/zip")
        else:
            data = save_as_pptx(final_pages)
            st.download_button("📥 PPTX 다운로드", data=data, file_name="slides.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

st.caption("※ 배경 이미지를 Fabric 객체로 삽입해 Streamlit 버전 호환 이슈를 우회했습니다.")
